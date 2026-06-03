"""
Generate a professional PowerPoint presentation for the OTT Streaming Platform project.
Demonstrates the architecture and cost savings using YouTube as storage via unlisted videos.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x0F, 0x1A)   # Deep dark background
BG_CARD      = RGBColor(0x1A, 0x1A, 0x2E)   # Card background
ACCENT       = RGBColor(0x6C, 0x5C, 0xE7)   # Purple accent
ACCENT_GOLD  = RGBColor(0xFF, 0xD7, 0x00)   # Gold accent
GREEN        = RGBColor(0x00, 0xB8, 0x94)   # Green (savings)
RED          = RGBColor(0xFF, 0x63, 0x48)   # Red (costs)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB0, 0xC0)
MID_GRAY     = RGBColor(0x80, 0x80, 0x90)
BLUE         = RGBColor(0x00, 0xB4, 0xD8)
ORANGE       = RGBColor(0xFF, 0xA0, 0x00)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        # Bullet
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
        # Remove existing bullets
        for old in pPr.findall(qn('a:buChar')):
            pPr.remove(old)
        for old in pPr.findall(qn('a:buNone')):
            pPr.remove(old)
        pPr.append(buChar)
        # Bullet color
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'[0:]})
        # Fix: get hex string
        hex_val = '%02X%02X%02X' % (bullet_color[0], bullet_color[1], bullet_color[2])
        srgbClr2 = buClr.makeelement(qn('a:srgbClr'), {'val': hex_val})
        buClr.append(srgbClr2)
        for old in pPr.findall(qn('a:buClr')):
            pPr.remove(old)
        pPr.append(buClr)

    return txBox


def add_icon_text(slide, left, top, icon, text, font_size=14, color=WHITE, icon_color=ACCENT):
    """Add an icon (emoji/symbol) followed by text."""
    txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{icon}  "
    run1.font.size = Pt(font_size + 2)
    run1.font.color.rgb = icon_color
    run1.font.name = "Segoe UI Emoji"
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(font_size)
    run2.font.color.rgb = color
    run2.font.name = "Calibri"
    return txBox


def add_stat_card(slide, left, top, width, height, number, label, color=ACCENT, bg=BG_CARD):
    """Add a stat card with a big number and label."""
    shape = add_shape(slide, left, top, width, height, bg, border_color=color, border_width=Pt(1.5))
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.6),
                 number, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.65), width - Inches(0.4), Inches(0.4),
                 label, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
#  CREATE PRESENTATION
# ═══════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # Blank layout

# ─────────────────────────────────────────────────────────────
# SLIDE 1: TITLE SLIDE
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)

# Decorative accent bar at top
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             "OTT Streaming Platform", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.8),
             "Zero-Cost Video Storage & Delivery Using YouTube's Unlisted Videos",
             font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Divider line
add_shape(slide, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.03), ACCENT)

# Description
add_text_box(slide, Inches(2), Inches(4.4), Inches(9), Inches(1),
             "A full-stack OTT platform built with Go, Flutter & PostgreSQL\nthat leverages YouTube as a free CDN for video hosting & delivery",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Tech stack badges
badge_y = Inches(5.8)
badges = [("Go", BLUE), ("Flutter", BLUE), ("PostgreSQL", GREEN), ("Docker", BLUE), ("YouTube API", RED)]
start_x = Inches(3.2)
for i, (label, clr) in enumerate(badges):
    x = start_x + Inches(i * 1.5)
    shape = add_shape(slide, x, badge_y, Inches(1.3), Inches(0.38), BG_CARD, border_color=clr, border_width=Pt(1))
    add_text_box(slide, x, badge_y + Inches(0.04), Inches(1.3), Inches(0.3),
                 label, font_size=11, color=clr, bold=True, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 2: THE PROBLEM
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), RED)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "⚠  The Problem: Video Hosting is Expensive", font_size=36, color=RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Traditional OTT platforms face massive infrastructure costs for video storage and delivery.",
             font_size=16, color=LIGHT_GRAY)

# Cost breakdown cards
costs = [
    ("$0.023/GB", "AWS S3 Storage\n(per month)", "💰"),
    ("$0.085/GB", "CloudFront CDN\nData Transfer", "🌐"),
    ("$0.024/min", "Video Transcoding\n(AWS MediaConvert)", "⚙"),
    ("~$5,000+", "Monthly Cost\nfor 10TB content", "📊"),
]
for i, (amount, desc, icon) in enumerate(costs):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.3)
    add_shape(slide, x, y, Inches(2.8), Inches(2.2), BG_CARD, border_color=RED, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.5),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.4), Inches(0.5),
                 amount, font_size=26, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.4), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom message
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.6),
             "For startups and indie developers, these costs can be a deal-breaker before the product even launches.",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key stat
add_shape(slide, Inches(3.5), Inches(5.8), Inches(6.3), Inches(1.0), BG_CARD, border_color=RED, border_width=Pt(1.5))
add_text_box(slide, Inches(3.7), Inches(5.9), Inches(5.9), Inches(0.5),
             "A 100-video library with 1080p content costs ~$500-$2,000/month on AWS",
             font_size=15, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.7), Inches(6.35), Inches(5.9), Inches(0.4),
             "Storage + CDN + Transcoding + Bandwidth",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 3: THE SOLUTION — YOUTUBE AS STORAGE
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "💡  The Solution: YouTube as Your Free CDN", font_size=36, color=GREEN, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Upload videos as Unlisted on YouTube → Embed them privately in your OTT app via YouTube IFrame API",
             font_size=17, color=LIGHT_GRAY)

# How it works - step by step
steps = [
    ("1", "Upload to YouTube", "Upload your video content\nas Unlisted — it won't\nappear in search or\nYouTube recommendations", "📤"),
    ("2", "Store Video ID", "The backend stores only\nthe YouTube Video ID\n(11 characters) in\nPostgreSQL database", "💾"),
    ("3", "Embed via API", "Flutter app uses YouTube\nIFrame Player API to\nembed and play videos\nwith custom controls", "▶"),
    ("4", "Full Control", "Custom UI hides YouTube\nbranding. Users see\nyour brand, not YouTube.\nYou control access.", "🎨"),
]

for i, (num, title, desc, icon) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.3)

    # Card
    add_shape(slide, x, y, Inches(2.8), Inches(3.0), BG_CARD, border_color=GREEN, border_width=Pt(1))

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), y + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    add_text_box(slide, x + Inches(1.05), y + Inches(0.22), Inches(0.5), Inches(0.45),
                 num, font_size=18, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Inches(0.15), y + Inches(0.85), Inches(2.5), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Inches(0.15), y + Inches(1.35), Inches(2.5), Inches(1.5),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom highlight
add_shape(slide, Inches(2.5), Inches(5.7), Inches(8.3), Inches(1.1), BG_CARD, border_color=GREEN, border_width=Pt(1.5))
add_text_box(slide, Inches(2.7), Inches(5.8), Inches(7.9), Inches(0.5),
             "YouTube handles: Storage • Transcoding • CDN • Adaptive Bitrate • Global Delivery",
             font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.7), Inches(6.3), Inches(7.9), Inches(0.4),
             "All for FREE — no API key required for embed playback",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 4: WHAT ARE UNLISTED VIDEOS?
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "🔒  YouTube Unlisted Videos — The Secret Weapon", font_size=36, color=ACCENT, bold=True)

# Three columns comparing video types
types = [
    ("Public", "❌ Not Suitable", RED,
     ["Appears in YouTube Search", "Shown in Recommendations", "Anyone can find it",
      "Indexed by Google", "Visible on your channel"]),
    ("Unlisted", "✅ Perfect for OTT", GREEN,
     ["NOT in YouTube Search", "NOT in Recommendations", "Only accessible via direct link",
      "NOT indexed by Google", "Embed works perfectly"]),
    ("Private", "❌ Too Restrictive", RED,
     ["Requires Google sign-in", "Max 50 viewers", "Cannot be embedded",
      "Not suitable for apps", "Breaks IFrame API"]),
]

for i, (vtype, verdict, clr, features) in enumerate(types):
    x = Inches(0.8 + i * 4.1)
    y = Inches(1.6)
    w = Inches(3.8)

    # Card
    add_shape(slide, x, y, w, Inches(5.2), BG_CARD, border_color=clr, border_width=Pt(1.5))

    # Type name
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.5),
                 vtype, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Verdict badge
    badge_shape = add_shape(slide, x + Inches(0.6), y + Inches(0.8), w - Inches(1.2), Inches(0.4), clr)
    add_text_box(slide, x + Inches(0.6), y + Inches(0.82), w - Inches(1.2), Inches(0.35),
                 verdict, font_size=12, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Features list
    for j, feat in enumerate(features):
        fy = y + Inches(1.5) + Inches(j * 0.6)
        prefix = "✓" if clr == GREEN else "✗"
        feat_color = GREEN if clr == GREEN else RGBColor(0x90, 0x90, 0xA0)
        add_text_box(slide, x + Inches(0.3), fy, w - Inches(0.5), Inches(0.5),
                     f"  {prefix}  {feat}", font_size=13, color=feat_color)


# ─────────────────────────────────────────────────────────────
# SLIDE 5: SYSTEM ARCHITECTURE
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🏗  System Architecture", font_size=36, color=BLUE, bold=True)

# Architecture components
components = [
    (Inches(0.5), Inches(1.8), Inches(3.0), Inches(4.5), "Flutter Mobile App", ACCENT,
     ["Home / Browse Screen", "Video Player (WebView)", "YouTube IFrame Player API", "Search & Categories",
      "Watchlist & History", "Premium Subscription", "User Authentication"]),
    (Inches(4.2), Inches(1.8), Inches(3.0), Inches(4.5), "Go Backend (Gin)", GREEN,
     ["REST API Server", "JWT Authentication", "Admin Middleware", "Video CRUD Handlers",
      "YouTube URL Parser", "User Management", "Payment Handler"]),
    (Inches(7.9), Inches(1.8), Inches(2.5), Inches(4.5), "PostgreSQL DB", ORANGE,
     ["Videos Table", "Users Table", "Categories Table", "Watchlist Table",
      "Watch History", "Payments Table"]),
    (Inches(11.0), Inches(1.8), Inches(2.0), Inches(4.5), "YouTube", RED,
     ["Video Storage", "Transcoding", "CDN Delivery", "Thumbnail CDN",
      "Adaptive Bitrate"]),
]

for x, y, w, h, title, clr, items in components:
    add_shape(slide, x, y, w, h, BG_CARD, border_color=clr, border_width=Pt(1.5))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), Inches(0.4),
                 title, font_size=15, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    # Divider
    add_shape(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), Inches(0.02), clr)
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.15), y + Inches(0.7 + j * 0.48), w - Inches(0.3), Inches(0.4),
                     f"• {item}", font_size=11, color=LIGHT_GRAY)

# Arrows between components
arrow_y = Inches(3.8)
for ax in [Inches(3.5), Inches(7.2), Inches(10.4)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, arrow_y, Inches(0.7), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MID_GRAY
    arrow.line.fill.background()

# Docker note
add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), BG_CARD, border_color=BLUE, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(6.55), Inches(12), Inches(0.6),
             "🐳  All services containerized with Docker Compose — single `docker-compose up` to deploy the entire stack",
             font_size=14, color=BLUE, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 6: COST COMPARISON
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "💰  Cost Comparison: Traditional vs. Our Approach", font_size=36, color=ACCENT_GOLD, bold=True)

# Table header
table_x = Inches(0.8)
table_y = Inches(1.5)
col_w = [Inches(3.5), Inches(4.0), Inches(4.0)]
row_h = Inches(0.6)

headers = ["Cost Category", "Traditional OTT (AWS)", "Our OTT (YouTube CDN)"]
header_colors = [ACCENT, RED, GREEN]
for i, (header, hclr) in enumerate(zip(headers, header_colors)):
    x = table_x + sum(w.inches for w in col_w[:i]) * 914400  # Convert back
    # Simpler approach
    cx = table_x
    for j in range(i):
        cx = Inches(cx.inches + col_w[j].inches)
    add_shape(slide, cx, table_y, col_w[i], row_h, BG_CARD, border_color=hclr, border_width=Pt(1))
    add_text_box(slide, cx + Inches(0.1), table_y + Inches(0.1), col_w[i] - Inches(0.2), Inches(0.4),
                 header, font_size=14, color=hclr, bold=True, alignment=PP_ALIGN.CENTER)

# Table rows
rows_data = [
    ("Video Storage (10TB)", "$230/month (S3)", "$0 (YouTube)"),
    ("CDN / Bandwidth (50TB)", "$4,250/month (CloudFront)", "$0 (YouTube CDN)"),
    ("Video Transcoding", "$500/month (MediaConvert)", "$0 (YouTube auto)"),
    ("Adaptive Bitrate", "$200/month (MediaPackage)", "$0 (YouTube auto)"),
    ("Thumbnail CDN", "$50/month", "$0 (YouTube CDN)"),
    ("Backend Server", "$50/month (EC2 / VPS)", "$5/month (VPS)"),
    ("Database", "$30/month (RDS)", "$5/month (Self-hosted)"),
    ("TOTAL (Monthly)", "~$5,310/month", "~$10/month"),
    ("TOTAL (Yearly)", "~$63,720/year", "~$120/year"),
]

for r, (cat, traditional, ours) in enumerate(rows_data):
    ry = Inches(table_y.inches + row_h.inches * (r + 1))
    row_bg = BG_CARD if r % 2 == 0 else RGBColor(0x15, 0x15, 0x25)
    is_total = r >= 7

    vals = [cat, traditional, ours]
    val_colors = [WHITE, RED if not is_total else RED, GREEN if not is_total else GREEN]
    for i, (val, vclr) in enumerate(zip(vals, val_colors)):
        cx = table_x
        for j in range(i):
            cx = Inches(cx.inches + col_w[j].inches)
        border = ACCENT_GOLD if is_total else None
        bw = Pt(2) if is_total else Pt(0)
        add_shape(slide, cx, ry, col_w[i], row_h, row_bg, border_color=border, border_width=bw)
        fs = 15 if is_total else 13
        b = is_total
        add_text_box(slide, cx + Inches(0.15), ry + Inches(0.12), col_w[i] - Inches(0.3), Inches(0.4),
                     val, font_size=fs, color=vclr, bold=b, alignment=PP_ALIGN.CENTER)

# Savings callout
add_shape(slide, Inches(3.5), Inches(6.7), Inches(6.3), Inches(0.6), GREEN, border_color=GREEN, border_width=Pt(1))
add_text_box(slide, Inches(3.5), Inches(6.72), Inches(6.3), Inches(0.55),
             "⬇ 99.8% Cost Reduction  |  Save ~$63,600/year",
             font_size=18, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 7: WHAT YOUTUBE GIVES YOU FOR FREE
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🎁  What YouTube Gives You for FREE", font_size=36, color=GREEN, bold=True)

# Feature cards grid (2 rows x 3 cols)
features = [
    ("Unlimited Storage", "Upload unlimited videos.\nNo storage limits.\nNo per-GB charges.", "♾"),
    ("Global CDN", "YouTube serves from 80+\ndata centers worldwide.\nFastest possible delivery.", "🌍"),
    ("Auto Transcoding", "YouTube automatically creates\n144p to 4K variants.\nNo encoding pipeline needed.", "⚡"),
    ("Adaptive Bitrate", "Viewer's quality auto-adjusts\nbased on their connection.\nSmooth playback always.", "📶"),
    ("Thumbnail CDN", "Auto-generated thumbnails\nat multiple resolutions.\nimg.youtube.com endpoint.", "🖼"),
    ("99.99% Uptime", "YouTube's infrastructure\nis one of the most reliable\nin the world.", "🛡"),
]

for i, (title, desc, icon) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.6 + row * 2.8)
    w = Inches(3.8)
    h = Inches(2.5)

    add_shape(slide, x, y, w, h, BG_CARD, border_color=GREEN, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.15), w - Inches(1), Inches(0.4),
                 title, font_size=17, color=GREEN, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.75), w - Inches(0.4), Inches(1.5),
                 desc, font_size=13, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────────────────
# SLIDE 8: HOW THE CODE WORKS
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "⚙  How It Works in Code", font_size=36, color=ACCENT, bold=True)

# Left: Backend flow
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5), BG_CARD, border_color=ACCENT, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.4),
             "Backend (Go) — Video Model", font_size=16, color=ACCENT, bold=True)

code_backend = """Video struct {
    YoutubeVideoID  string   // Only stores 11-char ID
    Title           string   // Video metadata
    ThumbnailURL    string   // Auto: youtube CDN
    IsPremium       bool     // Access control
    CategoryID      uint     // Organization
}

Admin → Pastes YouTube URL
→ Backend extracts Video ID (regex)
→ Stores ID + metadata in PostgreSQL
→ Thumbnail auto-generated from YouTube"""

add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.5), Inches(4.5),
             code_backend, font_size=12, color=LIGHT_GRAY, font_name="Consolas")

# Right: Frontend flow
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), BG_CARD, border_color=GREEN, border_width=Pt(1))
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
             "Frontend (Flutter) — Video Player", font_size=16, color=GREEN, bold=True)

code_frontend = """InAppWebView loads YouTube IFrame API:

player = new YT.Player('player', {
    videoId: '<youtube_video_id>',
    playerVars: {
        autoplay: 1,
        controls: 0,        // Hide YT controls
        modestbranding: 1,   // Minimal branding
        rel: 0,             // No related vids
        showinfo: 0         // No video info
    }
});

CSS hides all YouTube branding
→ User sees YOUR custom player UI"""

add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5),
             code_frontend, font_size=12, color=LIGHT_GRAY, font_name="Consolas")


# ─────────────────────────────────────────────────────────────
# SLIDE 9: KEY FEATURES
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "✨  Platform Features", font_size=36, color=ACCENT, bold=True)

# Feature grid
platform_features = [
    ("🔐", "User Authentication", "JWT-based login & signup\nSecure password hashing\nSession management", ACCENT),
    ("📺", "Custom Video Player", "YouTube embed with\ncustom controls overlay\nSeek bar & fullscreen", GREEN),
    ("👑", "Premium Tier System", "Free & premium content\nSubscription management\nContent gating", ACCENT_GOLD),
    ("📂", "Content Categories", "Organize by genre\nCategory-based browsing\nAdmin management", BLUE),
    ("📋", "Watchlist", "Save videos for later\nPersistent across sessions\nEasy add/remove", ORANGE),
    ("⏱", "Continue Watching", "Track watch progress\nResume from last position\nWatch history", GREEN),
    ("🔍", "Search", "Full-text search\nFilter by category\nReal-time results", BLUE),
    ("🛠", "Admin Dashboard", "Web-based admin panel\nCRUD for all entities\nUser & video management", RED),
]

for i, (icon, title, desc, clr) in enumerate(platform_features):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.5 + row * 2.8)
    w = Inches(2.95)
    h = Inches(2.5)

    add_shape(slide, x, y, w, h, BG_CARD, border_color=clr, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=26, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.65), w - Inches(0.3), Inches(0.4),
                 title, font_size=15, color=clr, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.1), w - Inches(0.3), Inches(1.2),
                 desc, font_size=12, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────────────────
# SLIDE 10: SCALING ANALYSIS
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "📈  Scaling Cost Analysis", font_size=36, color=ACCENT_GOLD, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "As you grow, the savings multiply exponentially",
             font_size=16, color=LIGHT_GRAY)

# Scale comparison table
scale_data = [
    ("Scale", "Videos", "Users", "Traditional Cost", "Our Cost", "You Save"),
    ("Hobby", "50", "100", "$1,500/mo", "$10/mo", "$1,490/mo"),
    ("Startup", "500", "5,000", "$8,000/mo", "$25/mo", "$7,975/mo"),
    ("Growth", "2,000", "50,000", "$25,000/mo", "$50/mo", "$24,950/mo"),
    ("Scale", "10,000", "500,000", "$80,000+/mo", "$100/mo", "$79,900/mo"),
]

col_widths = [Inches(1.6), Inches(1.5), Inches(1.6), Inches(2.4), Inches(2.0), Inches(2.6)]
for r, row_data in enumerate(scale_data):
    ry = Inches(1.9 + r * 0.65)
    for c, val in enumerate(row_data):
        cx = Inches(0.8)
        for j in range(c):
            cx = Inches(cx.inches + col_widths[j].inches)

        if r == 0:
            bg = ACCENT_GOLD
            txt_color = BG_DARK
            is_bold = True
        else:
            bg = BG_CARD if r % 2 == 0 else RGBColor(0x15, 0x15, 0x25)
            if c == 3:
                txt_color = RED
            elif c == 4:
                txt_color = GREEN
            elif c == 5:
                txt_color = ACCENT_GOLD
            else:
                txt_color = WHITE
            is_bold = c >= 3

        add_shape(slide, cx, ry, col_widths[c], Inches(0.55), bg)
        add_text_box(slide, cx + Inches(0.05), ry + Inches(0.08), col_widths[c] - Inches(0.1), Inches(0.4),
                     val, font_size=13, color=txt_color, bold=is_bold, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_shape(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.7), BG_CARD, border_color=ACCENT_GOLD, border_width=Pt(1.5))
add_text_box(slide, Inches(1.7), Inches(5.6), Inches(9.9), Inches(0.5),
             "Key Insight: Video costs scale linearly with traditional hosting.",
             font_size=16, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.7), Inches(6.1), Inches(9.9), Inches(0.9),
             "With YouTube CDN, your only scaling costs are backend compute and database,\n"
             "which grow at a fraction of the rate. A 10,000-video platform costs nearly\n"
             "the same as a 50-video platform in our architecture.",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 11: TECH STACK DEEP DIVE
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🛠  Tech Stack Deep Dive", font_size=36, color=BLUE, bold=True)

# Tech cards
tech_items = [
    ("Go + Gin", "High-performance backend\nCompiled to single binary\nFast HTTP routing\nMinimal memory footprint", BLUE,
     "Backend"),
    ("Flutter", "Cross-platform mobile app\nAndroid & iOS from one codebase\nMaterial Design 3\nSmooth 60fps UI", GREEN,
     "Mobile App"),
    ("PostgreSQL", "ACID-compliant database\nFull-text search built-in\nGORM as ORM\nDocker volume persistence", ORANGE,
     "Database"),
    ("Docker", "One-command deployment\nService isolation\nReproducible environments\nCompose orchestration", ACCENT,
     "DevOps"),
    ("YouTube API", "IFrame Player API\nNo API key for embeds\nJavaScript bridge to Flutter\nCustom player controls", RED,
     "Video Delivery"),
    ("JWT Auth", "Stateless authentication\nAdmin role middleware\nSecure token signing\nPassword hashing (bcrypt)", ACCENT_GOLD,
     "Security"),
]

for i, (title, desc, clr, category) in enumerate(tech_items):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.2)
    y = Inches(1.4 + row * 3.0)
    w = Inches(3.9)
    h = Inches(2.7)

    add_shape(slide, x, y, w, h, BG_CARD, border_color=clr, border_width=Pt(1))

    # Category badge
    badge_w = Inches(1.5)
    add_shape(slide, x + Inches(0.15), y + Inches(0.15), badge_w, Inches(0.3), clr)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), badge_w, Inches(0.28),
                 category, font_size=9, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.05), w - Inches(0.3), Inches(1.5),
                 desc, font_size=12, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────────────────
# SLIDE 12: TRADEOFFS & CONSIDERATIONS
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "⚖  Trade-offs & Considerations", font_size=36, color=ORANGE, bold=True)

# Two columns: Limitations vs Mitigations
# Left: Limitations
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), BG_CARD, border_color=ORANGE, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.4), Inches(0.4),
             "⚠  Potential Limitations", font_size=18, color=ORANGE, bold=True)

limitations = [
    "YouTube ToS — Ensure compliance with embedding terms",
    "No DRM — Content is accessible via YouTube if ID is known",
    "YouTube dependency — Platform risk if policies change",
    "Upload limits — 256GB or 12 hours per video max",
    "No offline downloads — Streaming only via YouTube",
]
for i, item in enumerate(limitations):
    add_text_box(slide, Inches(0.9), Inches(2.3 + i * 0.85), Inches(5.2), Inches(0.7),
                 f"▸ {item}", font_size=13, color=LIGHT_GRAY)

# Right: Mitigations
add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.5), BG_CARD, border_color=GREEN, border_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4),
             "✅  Mitigations & Strengths", font_size=18, color=GREEN, bold=True)

mitigations = [
    "Unlisted = not discoverable, only via your app",
    "YouTube explicitly supports embedding via API",
    "Modular design — can swap to S3 when revenue allows",
    "256GB handles 99%+ of content needs",
    "Perfect for MVP / early-stage / edu platforms",
]
for i, item in enumerate(mitigations):
    add_text_box(slide, Inches(7.3), Inches(2.3 + i * 0.85), Inches(5.2), Inches(0.7),
                 f"▸ {item}", font_size=13, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────────────────
# SLIDE 13: SUMMARY & CONCLUSION
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "🎯  Summary", font_size=36, color=GREEN, bold=True)

# Big stat cards
stats = [
    ("99.8%", "Cost\nReduction", GREEN),
    ("$0", "Video Storage\n& CDN Cost", ACCENT_GOLD),
    ("~$10/mo", "Total\nInfra Cost", BLUE),
    ("∞", "Scalable\nVideo Library", ACCENT),
]

for i, (num, label, clr) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    y = Inches(1.6)
    w = Inches(2.8)
    h = Inches(2.0)
    add_shape(slide, x, y, w, h, BG_CARD, border_color=clr, border_width=Pt(2))
    add_text_box(slide, x, y + Inches(0.2), w, Inches(0.7),
                 num, font_size=36, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.1), w, Inches(0.7),
                 label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key takeaways
takeaways = [
    "YouTube's unlisted video feature provides enterprise-grade CDN for free",
    "Our architecture eliminates the single largest cost for any OTT platform",
    "Full-featured platform: Auth, Subscriptions, Watchlist, History, Admin Panel",
    "Dockerized deployment — production-ready with a single command",
    "Perfect for startups, educators, and indie content platforms",
]

add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.2), BG_CARD, border_color=GREEN, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.4),
             "Key Takeaways", font_size=18, color=GREEN, bold=True)

for i, item in enumerate(takeaways):
    add_text_box(slide, Inches(0.9), Inches(4.6 + i * 0.48), Inches(11.5), Inches(0.4),
                 f"✓  {item}", font_size=14, color=LIGHT_GRAY)


# ─────────────────────────────────────────────────────────────
# SLIDE 14: THANK YOU
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
             "Thank You!", font_size=52, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.7),
             "OTT Streaming Platform — Built Smart, Scaled Free",
             font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.03), ACCENT)

add_text_box(slide, Inches(2), Inches(4.7), Inches(9), Inches(0.6),
             "Go • Flutter • PostgreSQL • Docker • YouTube IFrame API",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.6), Inches(9), Inches(0.6),
             "Questions? Let's discuss!",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "OTT_Platform_Presentation_v2.pptx")
prs.save(output_path)
print(f"[OK] Presentation saved to: {output_path}")
print(f"     Total slides: {len(prs.slides)}")
